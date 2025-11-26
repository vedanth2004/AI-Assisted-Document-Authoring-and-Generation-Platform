from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import FileResponse
from sqlalchemy.orm import Session
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.enum.style import WD_STYLE_TYPE
from pptx import Presentation
from pptx.util import Pt as PPTXPt, Inches as PPTXInches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as PPTXRGBColor
from io import BytesIO
import os
import re
from datetime import datetime

from database import get_db, User, Project, DocumentSection
from auth import get_current_user

router = APIRouter()

def format_paragraph_text(text):
    """Clean and format text, preserving structure"""
    # Remove excessive whitespace
    text = re.sub(r'\s+', ' ', text)
    # Preserve intentional line breaks for lists
    text = re.sub(r'\n\s*[-•*]\s*', '\n• ', text)
    return text.strip()

def add_formatted_content_to_paragraph(para, text, is_bullet=False):
    """Add formatted content to a paragraph with proper styling"""
    if is_bullet:
        para.style = 'List Bullet'
    else:
        para.style = 'Normal'
    
    # Split by bullet points or line breaks
    if '•' in text or '- ' in text or '* ' in text:
        # Handle bullet points
        lines = re.split(r'\n\s*[-•*]\s*', text)
        for i, line in enumerate(lines):
            if line.strip():
                if i > 0:
                    para = para._element.getparent().add_p('• ' + line.strip())
                else:
                    para.text = '• ' + line.strip()
    else:
        para.text = text
    
    # Formatting
    para_format = para.paragraph_format
    para_format.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
    para_format.line_spacing = 1.15
    para_format.space_after = Pt(6)
    
    # Font formatting
    for run in para.runs:
        run.font.size = Pt(11)
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(33, 33, 33)

@router.get("/{project_id}/download")
async def export_document(
    project_id: int,
    current_user: User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    project = db.query(Project).filter(
        Project.id == project_id,
        Project.user_id == current_user.id
    ).first()
    
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    
    db_sections = sorted(
        db.query(DocumentSection).filter(DocumentSection.project_id == project.id).all(),
        key=lambda x: x.section_index
    )
    
    if not db_sections:
        raise HTTPException(status_code=400, detail="No content to export")
    
    # Create temp directory if it doesn't exist
    temp_dir = "temp_exports"
    os.makedirs(temp_dir, exist_ok=True)
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{project.title.replace(' ', '_')}_{timestamp}"
    
    if project.document_type == "docx":
        # Create Word document with professional formatting
        doc = Document()
        
        # Set document margins
        doc_sections = doc.sections
        for section in doc_sections:
            section.top_margin = Inches(1)
            section.bottom_margin = Inches(1)
            section.left_margin = Inches(1)
            section.right_margin = Inches(1)
        
        # Add title page
        title_para = doc.add_paragraph()
        title_run = title_para.add_run(project.title if project.title != "Untitled Project" else project.topic)
        title_run.font.size = Pt(28)
        title_run.font.bold = True
        title_run.font.name = 'Calibri Light'
        title_run.font.color.rgb = RGBColor(31, 78, 121)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        title_para_format = title_para.paragraph_format
        title_para_format.space_after = Pt(12)
        
        # Add subtitle (topic)
        if project.title != "Untitled Project":
            subtitle_para = doc.add_paragraph()
            subtitle_run = subtitle_para.add_run(project.topic)
            subtitle_run.font.size = Pt(14)
            subtitle_run.font.italic = True
            subtitle_run.font.name = 'Calibri'
            subtitle_run.font.color.rgb = RGBColor(100, 100, 100)
            subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            subtitle_para_format = subtitle_para.paragraph_format
            subtitle_para_format.space_after = Pt(24)
        
        # Add date
        date_para = doc.add_paragraph()
        date_run = date_para.add_run(f"Generated on {datetime.now().strftime('%B %d, %Y')}")
        date_run.font.size = Pt(10)
        date_run.font.name = 'Calibri'
        date_run.font.color.rgb = RGBColor(128, 128, 128)
        date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        date_para_format = date_para.paragraph_format
        date_para_format.space_after = Pt(24)
        
        # Page break after title
        doc.add_page_break()
        
        # Add table of contents header
        toc_para = doc.add_heading('Table of Contents', level=1)
        toc_para_format = toc_para.paragraph_format
        toc_para_format.space_before = Pt(0)
        toc_para_format.space_after = Pt(12)
        
        # Add table of contents
        for idx, db_section in enumerate(db_sections, 1):
            toc_item = doc.add_paragraph(f"{idx}. {db_section.title}", style='List Number')
            toc_item_format = toc_item.paragraph_format
            toc_item_format.left_indent = Inches(0.25)
            toc_item_format.space_after = Pt(6)
        
        doc.add_page_break()
        
        # Add content sections
        for idx, db_section in enumerate(db_sections, 1):
            # Section heading
            heading = doc.add_heading(f"{idx}. {db_section.title}", level=1)
            heading_format = heading.paragraph_format
            heading_format.space_before = Pt(18)
            heading_format.space_after = Pt(12)
            
            # Format heading runs
            for run in heading.runs:
                run.font.size = Pt(18)
                run.font.bold = True
                run.font.name = 'Calibri'
                run.font.color.rgb = RGBColor(31, 78, 121)
            
            # Add content
            if db_section.content:
                content_text = format_paragraph_text(db_section.content)
                
                # Check if content has bullet points
                has_bullets = '•' in content_text or re.search(r'^\s*[-*]\s', content_text, re.MULTILINE)
                
                if has_bullets:
                    # Handle bullet points
                    lines = re.split(r'\n\s*[-•*]\s*', content_text)
                    for line in lines:
                        line = line.strip()
                        if line:
                            # Remove leading bullet if present
                            line = re.sub(r'^[-•*]\s*', '', line)
                            para = doc.add_paragraph(line, style='List Bullet')
                            para_format = para.paragraph_format
                            para_format.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
                            para_format.line_spacing = 1.15
                            para_format.left_indent = Inches(0.5)
                            para_format.space_after = Pt(6)
                            
                            for run in para.runs:
                                run.font.size = Pt(11)
                                run.font.name = 'Calibri'
                else:
                    # Regular paragraphs
                    paragraphs = content_text.split('\n\n')
                    for para_text in paragraphs:
                        para_text = para_text.strip()
                        if para_text:
                            para = doc.add_paragraph(para_text)
                            para_format = para.paragraph_format
                            para_format.line_spacing_rule = WD_LINE_SPACING.MULTIPLE
                            para_format.line_spacing = 1.15
                            para_format.first_line_indent = Inches(0.25)
                            para_format.space_after = Pt(12)
                            
                            for run in para.runs:
                                run.font.size = Pt(11)
                                run.font.name = 'Calibri'
                                run.font.color.rgb = RGBColor(33, 33, 33)
            
            # Add spacing between sections
            if idx < len(db_sections):
                doc.add_paragraph()
        
        # Add footer with page numbers
        section = doc.sections[0]
        footer = section.footer
        footer_para = footer.paragraphs[0]
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer_run = footer_para.add_run()
        footer_run.font.size = Pt(9)
        footer_run.font.name = 'Calibri'
        footer_run.font.color.rgb = RGBColor(128, 128, 128)
        # Note: python-docx doesn't support page numbers directly, but this structure is ready
        
        filepath = os.path.join(temp_dir, f"{filename}.docx")
        doc.save(filepath)
        
        return FileResponse(
            filepath,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=f"{filename}.docx"
        )
    
    else:  # pptx
        # Create PowerPoint presentation with professional formatting
        prs = Presentation()
        prs.slide_width = PPTXInches(10)
        prs.slide_height = PPTXInches(7.5)
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = project.title if project.title != "Untitled Project" else project.topic
        title_frame = title_shape.text_frame
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Format title
        for paragraph in title_frame.paragraphs:
            paragraph.font.size = PPTXPt(44)
            paragraph.font.bold = True
            paragraph.font.name = 'Calibri'
            paragraph.font.color.rgb = PPTXRGBColor(31, 78, 121)
            paragraph.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = project.topic if project.title != "Untitled Project" else f"Generated on {datetime.now().strftime('%B %d, %Y')}"
            subtitle_frame = subtitle.text_frame
            
            for paragraph in subtitle_frame.paragraphs:
                paragraph.font.size = PPTXPt(18)
                paragraph.font.name = 'Calibri'
                paragraph.font.color.rgb = PPTXRGBColor(100, 100, 100)
                paragraph.alignment = PP_ALIGN.CENTER
        
        # Content slides
        content_layout = prs.slide_layouts[1]
        
        for idx, db_section in enumerate(db_sections, 1):
            slide = prs.slides.add_slide(content_layout)
            
            # Slide title
            title_shape = slide.shapes.title
            title_shape.text = f"Slide {idx}: {db_section.title}"
            title_frame = title_shape.text_frame
            
            for paragraph in title_frame.paragraphs:
                paragraph.font.size = PPTXPt(32)
                paragraph.font.bold = True
                paragraph.font.name = 'Calibri'
                paragraph.font.color.rgb = PPTXRGBColor(31, 78, 121)
            
            # Content
            if len(slide.placeholders) > 1:
                content_shape = slide.placeholders[1]
                text_frame = content_shape.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = PPTXInches(0.5)
                text_frame.margin_right = PPTXInches(0.5)
                text_frame.margin_top = PPTXInches(0.5)
                text_frame.margin_bottom = PPTXInches(0.5)
                
                if db_section.content:
                    content_text = format_paragraph_text(db_section.content)
                    
                    # Check for bullet points
                    has_bullets = '•' in content_text or re.search(r'^\s*[-*]\s', content_text, re.MULTILINE)
                    
                    lines = re.split(r'\n+', content_text)
                    first_line = True
                    
                    for line in lines:
                        line = line.strip()
                        if not line:
                            continue
                        
                        # Remove bullet markers
                        line = re.sub(r'^[-•*]\s*', '', line)
                        
                        if first_line:
                            p = text_frame.paragraphs[0]
                            p.text = line
                            p.font.size = PPTXPt(18)
                            p.font.name = 'Calibri'
                            p.level = 0
                            if has_bullets:
                                p.font.bold = True
                            first_line = False
                        else:
                            p = text_frame.add_paragraph()
                            p.text = line
                            p.font.size = PPTXPt(16) if has_bullets else PPTXPt(18)
                            p.font.name = 'Calibri'
                            p.level = 0
                            if has_bullets:
                                p.font.bold = True
        
        filepath = os.path.join(temp_dir, f"{filename}.pptx")
        prs.save(filepath)
        
        return FileResponse(
            filepath,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"{filename}.pptx"
        )
